from pathlib import Path
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image

Image.MAX_IMAGE_PIXELS = None

FIGURE_DIR = Path("reports/figures")
TEMPLATE_PATH = Path("related-scripts-output/2025-10-08 Data and Coordination Working Group Update.pptx")
SOURCE_TEXT = "Source: EV Transition Forecasts"
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".bmp", ".gif"}


def output_path() -> Path:
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    return Path(f"reports/figures_summary_{timestamp}.pptx")


def format_title(file_path: Path) -> str:
    name = file_path.stem.replace("_", " ")
    return name.title()


def add_footer(slide, text: str) -> None:
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.4))
    tf = textbox.text_frame
    tf.text = text
    para = tf.paragraphs[0]
    para.font.size = Pt(12)
    para.font.italic = True


def main() -> None:
    images = sorted([
        p for p in FIGURE_DIR.glob("**/*")
        if p.suffix.lower() in IMAGE_EXTENSIONS and p.name != ".gitkeep"
    ], key=lambda p: p.as_posix())

    if not images:
        print("No figures found to include in presentation.")
        return

    template = TEMPLATE_PATH if TEMPLATE_PATH.exists() else None
    prs = Presentation(str(template)) if template else Presentation()

    title_slide_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_slide_layout)
    if title_slide.shapes.title:
        title_slide.shapes.title.text = "EV Transition Figures"
    if len(title_slide.placeholders) > 1:
        try:
            title_slide.placeholders[1].text = "Automatically generated from reports/figures"
        except IndexError:
            pass

    layout_candidates = [layout for layout in prs.slide_layouts if layout.name.lower() in {"title only", "blank"}]
    content_layout = layout_candidates[0] if layout_candidates else prs.slide_layouts[5]

    for image_path in images:
        slide = prs.slides.add_slide(content_layout)

        if slide.shapes.title:
            slide.shapes.title.text = format_title(image_path)
        else:
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
            title_box.text_frame.text = format_title(image_path)

        pic = slide.shapes.add_picture(str(image_path), Inches(0.5), Inches(1.2))

        max_width = Inches(9.0)
        max_height = Inches(5.5)

        if pic.width > max_width:
            ratio = max_width / pic.width
            pic.width = int(pic.width * ratio)
            pic.height = int(pic.height * ratio)
        if pic.height > max_height:
            ratio = max_height / pic.height
            pic.width = int(pic.width * ratio)
            pic.height = int(pic.height * ratio)

        pic.left = int((prs.slide_width - pic.width) / 2)
        pic.top = Inches(1.4)

        add_footer(slide, SOURCE_TEXT)

    output = output_path()
    prs.save(output)
    print(f"Presentation saved to {output}")


if __name__ == "__main__":
    main()
